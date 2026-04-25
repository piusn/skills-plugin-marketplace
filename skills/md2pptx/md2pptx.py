#!/usr/bin/env python3
"""
md2pptx — Convert Markdown files to *presentation-quality* PowerPoint decks.

Unlike a raw markdown-to-slides dump, this tool **plans** the presentation:

  • Key talking points go ON the slide (concise, visual).
  • Details, code snippets, full tables, and extended explanations go
    into the **Speaker Notes** so the presenter can reference them.
  • Sections are consolidated to keep the slide count tight.

Usage:
    python md2pptx.py <input.md> [output.pptx] [--title "Title"]

Slide-planning rules:
    # H1                → Title slide
    ## H2               → Section divider slide
    ### H3              → Content slide title
    Bullets (≤ MAX)     → On slide; overflow → speaker notes
    Tables (≤ MAX rows) → On slide; overflow → speaker notes
    Code blocks         → Always speaker notes (summary on slide)
    Long paragraphs     → First sentence on slide; rest → notes
    Blockquotes         → Callout on slide
"""

import argparse
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# ---------------------------------------------------------------------------
# Presentation limits — what fits on a single slide
# ---------------------------------------------------------------------------
MAX_BULLETS_ON_SLIDE = 6
MAX_TABLE_ROWS_ON_SLIDE = 5
MAX_NUMBERED_ON_SLIDE = 6
MAX_PARAGRAPH_CHARS = 200       # characters before spilling to notes


# ---------------------------------------------------------------------------
# Theme / colours
# ---------------------------------------------------------------------------
class Theme:
    BG_TITLE = RGBColor(0x0F, 0x1B, 0x2D)
    BG_SECTION = RGBColor(0x0F, 0x1B, 0x2D)
    BG_CONTENT = RGBColor(0xFF, 0xFF, 0xFF)

    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    BLACK = RGBColor(0x1A, 0x1A, 0x2E)
    GREY = RGBColor(0x6B, 0x70, 0x80)
    ACCENT = RGBColor(0x00, 0x78, 0xD4)
    ACCENT_LIGHT = RGBColor(0xDE, 0xEC, 0xF9)
    WARNING = RGBColor(0xE7, 0x4C, 0x3C)
    SUCCESS = RGBColor(0x2E, 0xCC, 0x71)

    TABLE_HEADER_BG = RGBColor(0x0F, 0x1B, 0x2D)
    TABLE_HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
    TABLE_ROW_ALT = RGBColor(0xF0, 0xF4, 0xF8)
    TABLE_BORDER = RGBColor(0xD0, 0xD5, 0xDD)

    CODE_BG = RGBColor(0xF6, 0xF8, 0xFA)
    CODE_FG = RGBColor(0x24, 0x29, 0x2E)

    FONT_HEADING = "Segoe UI Semibold"
    FONT_BODY = "Segoe UI"
    FONT_CODE = "Cascadia Code"


# ---------------------------------------------------------------------------
# Markdown parser  (unchanged — produces structured blocks)
# ---------------------------------------------------------------------------
class MarkdownParser:
    """Parse markdown into a list of structured blocks."""

    def parse(self, text: str) -> list[dict]:
        lines = text.split("\n")
        blocks: list[dict] = []
        i = 0

        while i < len(lines):
            line = lines[i]

            if re.match(r"^-{3,}\s*$", line) or re.match(r"^\*{3,}\s*$", line):
                i += 1
                continue

            if line.strip().startswith("```"):
                lang = line.strip().lstrip("`").strip()
                code_lines: list[str] = []
                i += 1
                while i < len(lines) and not lines[i].strip().startswith("```"):
                    code_lines.append(lines[i])
                    i += 1
                i += 1
                blocks.append({"type": "code", "lang": lang, "content": "\n".join(code_lines)})
                continue

            m = re.match(r"^(#{1,6})\s+(.+)$", line)
            if m:
                blocks.append({"type": "heading", "level": len(m.group(1)), "content": m.group(2).strip()})
                i += 1
                continue

            if "|" in line and i + 1 < len(lines) and re.match(r"^\s*\|[\s\-:|]+\|\s*$", lines[i + 1]):
                table_lines: list[str] = []
                while i < len(lines) and "|" in lines[i]:
                    stripped = lines[i].strip()
                    if re.match(r"^\|[\s\-:|]+\|$", stripped):
                        i += 1
                        continue
                    table_lines.append(stripped)
                    i += 1
                rows = [[c.strip() for c in tl.strip("|").split("|")] for tl in table_lines]
                if rows:
                    blocks.append({"type": "table", "header": rows[0], "rows": rows[1:]})
                continue

            if line.strip().startswith(">"):
                quote_lines: list[str] = []
                while i < len(lines) and lines[i].strip().startswith(">"):
                    quote_lines.append(re.sub(r"^>\s*", "", lines[i].strip()))
                    i += 1
                blocks.append({"type": "blockquote", "content": " ".join(quote_lines)})
                continue

            m_bullet = re.match(r"^(\s*)[-*]\s+(\[[ x]\]\s+)?(.+)$", line)
            if m_bullet:
                items = [self._parse_bullet(m_bullet)]
                i += 1
                while i < len(lines):
                    m2 = re.match(r"^(\s*)[-*]\s+(\[[ x]\]\s+)?(.+)$", lines[i])
                    if m2:
                        items.append(self._parse_bullet(m2))
                        i += 1
                    else:
                        break
                blocks.append({"type": "bullets", "items": items})
                continue

            m_num = re.match(r"^\s*\d+\.\s+(.+)$", line)
            if m_num:
                items = [{"text": m_num.group(1).strip(), "level": 0}]
                i += 1
                while i < len(lines):
                    m2 = re.match(r"^\s*\d+\.\s+(.+)$", lines[i])
                    if m2:
                        items.append({"text": m2.group(1).strip(), "level": 0})
                        i += 1
                    else:
                        break
                blocks.append({"type": "numbered", "items": items})
                continue

            if not line.strip():
                i += 1
                continue

            para_lines = [line.strip()]
            i += 1
            while (i < len(lines) and lines[i].strip()
                   and not lines[i].startswith("#") and not lines[i].startswith("|")
                   and not lines[i].startswith(">") and not lines[i].startswith("```")
                   and not lines[i].strip().startswith("- ")
                   and not re.match(r"^\s*\d+\.\s+", lines[i])
                   and not re.match(r"^-{3,}\s*$", lines[i])):
                para_lines.append(lines[i].strip())
                i += 1
            blocks.append({"type": "paragraph", "content": " ".join(para_lines)})

        return blocks

    @staticmethod
    def _parse_bullet(m) -> dict:
        return {
            "text": m.group(3).strip(),
            "level": len(m.group(1)) // 2,
            "checkbox": m.group(2) is not None,
            "checked": bool(m.group(2) and "x" in m.group(2)),
        }


# ---------------------------------------------------------------------------
# Presentation planner
# ---------------------------------------------------------------------------
class PresentationPlanner:
    """Analyse parsed blocks and produce a *slide plan*.

    Each planned slide is a dict:
        {
            "kind": "title" | "section" | "content",
            "title": str,
            "subtitle": str,          # title/section only
            "slide_blocks": [block],   # blocks rendered ON the slide
            "notes_blocks": [block],   # blocks rendered in speaker notes
        }
    """

    def plan(self, blocks: list[dict], title: str | None = None) -> list[dict]:
        slides: list[dict] = []

        doc_title = title or "Presentation"
        doc_subtitle = ""
        idx = 0

        # First H1 → title
        if blocks and blocks[0]["type"] == "heading" and blocks[0]["level"] == 1:
            doc_title = blocks[0]["content"]
            idx = 1
        if idx < len(blocks) and blocks[idx]["type"] == "paragraph":
            doc_subtitle = blocks[idx]["content"]
            idx += 1

        slides.append({"kind": "title", "title": doc_title, "subtitle": doc_subtitle,
                        "slide_blocks": [], "notes_blocks": []})

        # Group remaining blocks into sections (H2) and sub-sections (H3)
        current_title: str | None = None
        current_slide: list[dict] = []
        current_notes: list[dict] = []

        def flush():
            nonlocal current_slide, current_notes
            if current_title and (current_slide or current_notes):
                # Enforce height budget: move overflow blocks to notes
                current_slide, current_notes = self._enforce_budget(
                    current_slide, current_notes)
                slides.append({
                    "kind": "content",
                    "title": current_title,
                    "subtitle": "",
                    "slide_blocks": current_slide,
                    "notes_blocks": current_notes,
                })
            current_slide = []
            current_notes = []

        for block in blocks[idx:]:
            if block["type"] == "heading":
                if block["level"] == 2:
                    flush()
                    current_title = block["content"]
                    slides.append({"kind": "section", "title": current_title,
                                   "subtitle": "", "slide_blocks": [], "notes_blocks": []})
                elif block["level"] == 3:
                    flush()
                    current_title = block["content"]
                elif block["level"] >= 4:
                    # Sub-headings: add as bold text on slide, treat as slide content
                    current_slide.append(block)
                continue

            # ── Decide: slide vs notes ──
            on_slide, in_notes = self._split_block(block)
            if on_slide:
                current_slide.append(on_slide)
            if in_notes:
                current_notes.append(in_notes)

        flush()

        # ── Consolidation pass ──
        # 1. Merge small adjacent content slides to reduce count.
        # 2. Eliminate section dividers that precede a single content slide.
        slides = self._consolidate(slides)
        slides = self._fold_lonely_sections(slides)

        # Closing slide
        slides.append({"kind": "title", "title": "Thank You",
                        "subtitle": "Questions & Discussion",
                        "slide_blocks": [], "notes_blocks": []})
        return slides

    # ── Consolidation ──────────────────────────────────────────────────────

    MAX_CONSOLIDATED_HEIGHT = 4.8  # inches of content before splitting

    def _consolidate(self, slides: list[dict]) -> list[dict]:
        """Merge small adjacent content slides into single slides."""
        result: list[dict] = []
        i = 0
        while i < len(slides):
            s = slides[i]
            if s["kind"] != "content":
                result.append(s)
                i += 1
                continue

            merged_blocks = list(s["slide_blocks"])
            merged_notes = list(s["notes_blocks"])
            merged_titles = [s["title"]]
            est = sum(_estimate_height(b) for b in merged_blocks)

            j = i + 1
            while j < len(slides) and slides[j]["kind"] == "content":
                next_s = slides[j]
                next_est = sum(_estimate_height(b) for b in next_s["slide_blocks"])
                if est + next_est + 0.45 <= self.MAX_CONSOLIDATED_HEIGHT:
                    merged_blocks.append({"type": "heading", "level": 4,
                                          "content": next_s["title"]})
                    merged_blocks.extend(next_s["slide_blocks"])
                    merged_notes.extend(next_s["notes_blocks"])
                    merged_titles.append(next_s["title"])
                    est += next_est + 0.45
                    j += 1
                else:
                    break

            if j > i + 1:
                result.append({
                    "kind": "content",
                    "title": merged_titles[0],
                    "subtitle": "",
                    "slide_blocks": merged_blocks,
                    "notes_blocks": merged_notes,
                })
            else:
                result.append(s)
            i = j

        return result

    @staticmethod
    def _fold_lonely_sections(slides: list[dict]) -> list[dict]:
        """Remove section dividers that have only one content slide after them.

        The section title becomes the content slide title instead, saving a slide.
        """
        result: list[dict] = []
        i = 0
        while i < len(slides):
            s = slides[i]
            if s["kind"] == "section" and i + 1 < len(slides):
                nxt = slides[i + 1]
                # Count how many content slides follow this section
                content_count = 0
                j = i + 1
                while j < len(slides) and slides[j]["kind"] == "content":
                    content_count += 1
                    j += 1
                if content_count == 1:
                    # Fold: skip section slide, use section title as content title
                    folded = dict(nxt)
                    folded["title"] = s["title"]
                    result.append(folded)
                    i += 2
                    continue
            result.append(s)
            i += 1
        return result

    # ── Per-block splitting logic ──────────────────────────────────────────

    SLIDE_BUDGET = 5.0  # max inches of content on a single slide

    @staticmethod
    def _enforce_budget(slide_blocks: list[dict],
                        notes_blocks: list[dict]) -> tuple[list[dict], list[dict]]:
        """Move excess on-slide blocks to notes when total height exceeds budget."""
        budget = PresentationPlanner.SLIDE_BUDGET
        kept: list[dict] = []
        overflow: list[dict] = list(notes_blocks)
        used = 0.0

        for b in slide_blocks:
            h = _estimate_height(b)
            if used + h <= budget:
                kept.append(b)
                used += h
            else:
                overflow.append(b)

        return kept, overflow

    @staticmethod
    def _split_block(block: dict) -> tuple[dict | None, dict | None]:
        """Return (on_slide_block_or_None, notes_block_or_None)."""
        btype = block["type"]

        # CODE → always in notes, with a short on-slide indicator
        if btype == "code":
            lang = block.get("lang", "")
            label = f"Code ({lang})" if lang else "Code"
            first_lines = block["content"].strip().split("\n")[:3]
            preview = "\n".join(first_lines)
            if len(block["content"].strip().split("\n")) > 3:
                preview += "\n..."
            summary = {"type": "blockquote", "content": f"📝 {label} — see speaker notes for full snippet"}
            return summary, block

        # BULLETS → keep first N on slide, overflow to notes
        if btype == "bullets":
            items = block["items"]
            if len(items) <= MAX_BULLETS_ON_SLIDE:
                return block, None
            on = {"type": "bullets", "items": items[:MAX_BULLETS_ON_SLIDE]}
            off = {"type": "bullets", "items": items[MAX_BULLETS_ON_SLIDE:]}
            return on, off

        # NUMBERED → same logic
        if btype == "numbered":
            items = block["items"]
            if len(items) <= MAX_NUMBERED_ON_SLIDE:
                return block, None
            on = {"type": "numbered", "items": items[:MAX_NUMBERED_ON_SLIDE]}
            off = {"type": "numbered", "items": items[MAX_NUMBERED_ON_SLIDE:]}
            return on, off

        # TABLE → keep header + first N rows on slide, rest to notes
        if btype == "table":
            rows = block["rows"]
            if len(rows) <= MAX_TABLE_ROWS_ON_SLIDE:
                return block, None
            on = {"type": "table", "header": block["header"],
                  "rows": rows[:MAX_TABLE_ROWS_ON_SLIDE]}
            off = {"type": "table", "header": block["header"],
                   "rows": rows[MAX_TABLE_ROWS_ON_SLIDE:]}
            return on, off

        # PARAGRAPH → if too long, first sentence on slide, rest in notes
        if btype == "paragraph":
            text = block["content"]
            if len(text) <= MAX_PARAGRAPH_CHARS:
                return block, None
            # Split on first sentence boundary
            m = re.match(r"^(.+?[.!?])\s+(.+)$", text)
            if m:
                return ({"type": "paragraph", "content": m.group(1)},
                        {"type": "paragraph", "content": text})
            return block, {"type": "paragraph", "content": text}

        # BLOCKQUOTE → always on slide (usually short callouts)
        if btype == "blockquote":
            return block, None

        # Fallback
        return block, None


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def _strip_md_inline(text: str) -> str:
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    text = re.sub(r"\[(.+?)\]\(.+?\)", r"\1", text)
    return text


def _add_formatted_text(paragraph, text: str, font_size=Pt(14), color=None, bold=None):
    if color is None:
        color = Theme.BLACK
    pattern = r"(\*\*.*?\*\*|`[^`]+`|\[.*?\]\(.*?\)|\*.*?\*)"
    segments = re.split(pattern, text)
    for seg in segments:
        if not seg:
            continue
        run = paragraph.add_run()
        run.font.size = font_size
        if seg.startswith("**") and seg.endswith("**"):
            run.text = seg[2:-2]
            run.font.bold = True
            run.font.color.rgb = color
        elif seg.startswith("`") and seg.endswith("`"):
            run.text = seg[1:-1]
            run.font.name = Theme.FONT_CODE
            run.font.color.rgb = Theme.ACCENT
            run.font.size = Pt(font_size.pt - 1) if hasattr(font_size, "pt") else Pt(13)
        elif seg.startswith("["):
            m = re.match(r"\[(.+?)\]\((.+?)\)", seg)
            if m:
                run.text = m.group(1)
                run.font.color.rgb = Theme.ACCENT
                run.font.underline = True
            else:
                run.text = seg
                run.font.color.rgb = color
        elif seg.startswith("*") and seg.endswith("*"):
            run.text = seg[1:-1]
            run.font.italic = True
            run.font.color.rgb = color
        else:
            run.text = seg
            run.font.color.rgb = color
            if bold is not None:
                run.font.bold = bold
        run.font.name = run.font.name or Theme.FONT_BODY


def _block_to_plain_text(block: dict) -> str:
    """Convert a parsed block back to readable plain text for speaker notes."""
    btype = block["type"]
    if btype == "paragraph":
        return _strip_md_inline(block["content"])
    if btype == "blockquote":
        return f"> {_strip_md_inline(block['content'])}"
    if btype == "code":
        lang = block.get("lang", "")
        return f"```{lang}\n{block['content']}\n```"
    if btype == "bullets":
        lines = []
        for item in block["items"]:
            indent = "  " * item["level"]
            prefix = ""
            if item.get("checkbox"):
                prefix = "[x] " if item.get("checked") else "[ ] "
            lines.append(f"{indent}• {prefix}{_strip_md_inline(item['text'])}")
        return "\n".join(lines)
    if btype == "numbered":
        return "\n".join(f"{i+1}. {_strip_md_inline(it['text'])}"
                         for i, it in enumerate(block["items"]))
    if btype == "table":
        hdr = " | ".join(_strip_md_inline(c) for c in block["header"])
        sep = " | ".join("---" for _ in block["header"])
        rows = "\n".join(" | ".join(_strip_md_inline(c) for c in r) for r in block["rows"])
        return f"{hdr}\n{sep}\n{rows}"
    if btype == "heading":
        return _strip_md_inline(block["content"])
    return str(block.get("content", ""))


# ---------------------------------------------------------------------------
# Slide builder  (placeholder-based for Designer compatibility)
# ---------------------------------------------------------------------------
class SlideBuilder:
    """Build PowerPoint slides using proper layouts so Designer can suggest designs.

    Layout mapping (default template):
        0 = Title Slide    (CENTER_TITLE + SUBTITLE)
        1 = Title & Content (TITLE + OBJECT)
        2 = Section Header  (TITLE + BODY)
        5 = Title Only       (TITLE)
        6 = Blank
    """

    # Layout indices in the default python-pptx template
    LAYOUT_TITLE = 0
    LAYOUT_CONTENT = 1
    LAYOUT_SECTION = 2
    LAYOUT_TITLE_ONLY = 5
    LAYOUT_BLANK = 6

    def __init__(self, prs: Presentation, title: str = ""):
        self.prs = prs
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.title = title
        self.slide_num = 0

    # ── helpers ──

    @staticmethod
    def _widen_placeholders(slide):
        """Stretch placeholders to fill 16:9 widescreen (13.333 in).

        The default python-pptx template is 10 in wide so placeholders
        are too narrow for widescreen.  Preserves all four dimensions
        to avoid the python-pptx inherited-placeholder reset quirk.
        """
        margin = Inches(0.6)
        usable = Inches(13.333) - 2 * margin

        for ph in slide.placeholders:
            idx = ph.placeholder_format.idx
            if idx >= 10:
                continue
            # Save all dimensions BEFORE modifying any
            orig_top = ph.top
            orig_height = ph.height
            # Set all four at once
            ph.left = margin
            ph.width = usable
            ph.top = orig_top
            ph.height = orig_height

    @staticmethod
    def _set_notes(slide, notes_blocks: list[dict]):
        """Write blocks into the slide's speaker notes."""
        if not notes_blocks:
            return
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.clear()
        first = True
        for block in notes_blocks:
            text = _block_to_plain_text(block)
            if not text.strip():
                continue
            if first:
                tf.paragraphs[0].text = text
                first = False
            else:
                p = tf.add_paragraph()
                p.text = text

    # ── slide types ──

    def add_title_slide(self, title: str, subtitle: str = "", notes=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.LAYOUT_TITLE])
        self._widen_placeholders(slide)
        self.slide_num += 1

        # Use the CENTER_TITLE placeholder (idx 0)
        title_ph = slide.placeholders[0]
        title_ph.text = _strip_md_inline(title)
        for p in title_ph.text_frame.paragraphs:
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.name = Theme.FONT_HEADING

        # Use the SUBTITLE placeholder (idx 1)
        if subtitle:
            sub_ph = slide.placeholders[1]
            sub_ph.text = _strip_md_inline(subtitle)
            for p in sub_ph.text_frame.paragraphs:
                p.font.size = Pt(18)
                p.font.name = Theme.FONT_BODY

        if notes:
            self._set_notes(slide, notes)
        return slide

    def add_section_slide(self, title: str, notes=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.LAYOUT_SECTION])
        self._widen_placeholders(slide)
        self.slide_num += 1

        # TITLE placeholder (idx 0)
        title_ph = slide.placeholders[0]
        title_ph.text = _strip_md_inline(title)
        for p in title_ph.text_frame.paragraphs:
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.name = Theme.FONT_HEADING

        if notes:
            self._set_notes(slide, notes)
        return slide

    def add_content_slide(self, title: str, slide_blocks: list[dict],
                          notes_blocks: list[dict] | None = None):
        # Separate text-based blocks from table/code blocks
        text_blocks = []
        shape_blocks = []
        for b in slide_blocks:
            if b["type"] in ("table", "code"):
                shape_blocks.append(b)
            else:
                text_blocks.append(b)

        # Choose layout based on content type
        if text_blocks and not shape_blocks:
            # Pure text → use Title & Content layout for full Designer support
            slide = self._build_text_content_slide(title, text_blocks)
        elif shape_blocks and not text_blocks:
            # Only tables/code → Title Only + shapes
            slide = self._build_shape_content_slide(title, shape_blocks)
        else:
            # Mixed → Title & Content for text, shapes added below
            slide = self._build_mixed_content_slide(title, text_blocks, shape_blocks)

        self.slide_num += 1
        if notes_blocks:
            self._set_notes(slide, notes_blocks)
        return slide

    # ── content slide builders ──

    def _build_text_content_slide(self, title: str, blocks: list[dict]):
        """Build slide using Layout 1 (Title & Content) with text in placeholder."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.LAYOUT_CONTENT])
        self._widen_placeholders(slide)

        # TITLE placeholder (idx 0)
        slide.placeholders[0].text = _strip_md_inline(title)
        for p in slide.placeholders[0].text_frame.paragraphs:
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.name = Theme.FONT_HEADING

        # OBJECT / CONTENT placeholder (idx 1) — add all text blocks here
        content_ph = slide.placeholders[1]
        tf = content_ph.text_frame
        tf.clear()
        self._fill_text_frame(tf, blocks)
        return slide

    def _build_shape_content_slide(self, title: str, blocks: list[dict]):
        """Build slide using Layout 5 (Title Only) + shape objects."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.LAYOUT_TITLE_ONLY])
        self._widen_placeholders(slide)

        slide.placeholders[0].text = _strip_md_inline(title)
        for p in slide.placeholders[0].text_frame.paragraphs:
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.name = Theme.FONT_HEADING

        y_pos = Inches(1.8)
        max_y = Inches(6.8)
        for block in blocks:
            if y_pos >= max_y:
                break
            y_pos = self._render_shape_block(slide, block, y_pos, max_y)
        return slide

    def _build_mixed_content_slide(self, title: str, text_blocks: list[dict],
                                   shape_blocks: list[dict]):
        """Build slide with text in placeholder and shapes below."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.LAYOUT_CONTENT])
        self._widen_placeholders(slide)

        slide.placeholders[0].text = _strip_md_inline(title)
        for p in slide.placeholders[0].text_frame.paragraphs:
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.name = Theme.FONT_HEADING

        content_ph = slide.placeholders[1]
        tf = content_ph.text_frame
        tf.clear()
        self._fill_text_frame(tf, text_blocks)

        # Shrink placeholder to make room for shapes, preserving width
        text_height = sum(_estimate_height(b) for b in text_blocks)
        ph_bottom = content_ph.top + Inches(min(text_height + 0.3, 3.5))
        orig_width = content_ph.width
        content_ph.height = int(ph_bottom - content_ph.top)
        content_ph.width = orig_width

        y_pos = ph_bottom + Inches(0.2)
        max_y = Inches(6.8)
        for block in shape_blocks:
            if y_pos >= max_y:
                break
            y_pos = self._render_shape_block(slide, block, y_pos, max_y)
        return slide

    # ── fill a text frame with blocks ──

    def _fill_text_frame(self, tf, blocks: list[dict]):
        """Add parsed blocks as paragraphs in a text frame (placeholder)."""
        first = True
        for block in blocks:
            btype = block["type"]

            if btype == "paragraph":
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.space_before = Pt(4)
                p.space_after = Pt(4)
                _add_formatted_text(p, block["content"], Pt(14))

            elif btype == "bullets":
                for item in block["items"]:
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    first = False
                    p.level = min(item["level"], 2)
                    p.space_before = Pt(2)
                    p.space_after = Pt(2)
                    prefix = ""
                    if item.get("checkbox"):
                        prefix = "☑ " if item.get("checked") else "☐ "
                    _add_formatted_text(p, f"{prefix}{item['text']}", Pt(13))

            elif btype == "numbered":
                for idx, item in enumerate(block["items"]):
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    first = False
                    p.space_before = Pt(2)
                    p.space_after = Pt(2)
                    _add_formatted_text(p, f"{idx + 1}. {item['text']}", Pt(13))

            elif btype == "blockquote":
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.space_before = Pt(6)
                p.space_after = Pt(6)
                _add_formatted_text(p, f"▎ {block['content']}", Pt(12), Theme.GREY)
                p.font.italic = True

            elif btype == "heading":
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.space_before = Pt(8)
                p.space_after = Pt(2)
                size = {3: Pt(20), 4: Pt(16), 5: Pt(14), 6: Pt(12)}.get(
                    block["level"], Pt(16))
                run = p.add_run()
                run.text = _strip_md_inline(block["content"])
                run.font.size = size
                run.font.bold = True
                run.font.color.rgb = Theme.ACCENT
                run.font.name = Theme.FONT_HEADING

    # ── render shapes (tables, code) that can't go in text placeholders ──

    def _render_shape_block(self, slide, block, y_pos, max_y):
        btype = block["type"]

        if btype == "table":
            rows_count = len(block["rows"]) + 1
            cols_count = len(block["header"])
            table_width = min(Inches(12.0), Inches(2.5) * cols_count)
            row_height = Inches(0.35)
            table_height = row_height * rows_count

            if y_pos + table_height > max_y:
                available = int((max_y - y_pos) / row_height) - 1
                if available <= 0:
                    return y_pos
                block = dict(block)
                block["rows"] = block["rows"][:available]
                rows_count = len(block["rows"]) + 1
                table_height = row_height * rows_count

            tbl_shape = slide.shapes.add_table(rows_count, cols_count,
                                               Inches(0.6), y_pos,
                                               table_width, table_height)
            tbl = tbl_shape.table
            for ci, cell_text in enumerate(block["header"]):
                cell = tbl.cell(0, ci)
                cell.text = _strip_md_inline(cell_text)
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(11)
                    p.font.bold = True
                    p.font.name = Theme.FONT_BODY
                    p.alignment = PP_ALIGN.LEFT
            for ri, row in enumerate(block["rows"]):
                for ci, cell_text in enumerate(row):
                    if ci >= cols_count:
                        break
                    cell = tbl.cell(ri + 1, ci)
                    cell.text = _strip_md_inline(cell_text)
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(10)
                        p.font.name = Theme.FONT_BODY
                        p.alignment = PP_ALIGN.LEFT
            return y_pos + table_height + Inches(0.2)

        if btype == "code":
            code_text = block["content"]
            code_lines = code_text.count("\n") + 1
            box_height = min(Inches(0.25) * code_lines + Inches(0.3), max_y - y_pos)
            bg_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                              Inches(0.6), y_pos, Inches(12.0), box_height)
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = Theme.CODE_BG
            bg_shape.line.color.rgb = Theme.TABLE_BORDER
            bg_shape.line.width = Pt(1)
            txBox = slide.shapes.add_textbox(Inches(0.8), y_pos + Inches(0.1),
                                             Inches(11.6), box_height - Inches(0.2))
            tf = txBox.text_frame
            tf.word_wrap = False
            run = tf.paragraphs[0].add_run()
            display_lines = code_text.split("\n")
            max_code_lines = max(1, int((box_height - Inches(0.2)) / Inches(0.2)))
            if len(display_lines) > max_code_lines:
                display_lines = display_lines[:max_code_lines] + ["  ..."]
            run.text = "\n".join(display_lines)
            run.font.size = Pt(10)
            run.font.name = Theme.FONT_CODE
            run.font.color.rgb = Theme.CODE_FG
            return y_pos + box_height + Inches(0.15)

        return y_pos


# ---------------------------------------------------------------------------
# Height estimator (for multi-slide splitting)
# ---------------------------------------------------------------------------
def _estimate_height(block: dict) -> float:
    btype = block["type"]
    if btype == "paragraph":
        return 0.5
    if btype in ("bullets", "numbered"):
        return 0.35 * len(block.get("items", []))
    if btype == "table":
        return 0.35 * (len(block.get("rows", [])) + 1) + 0.2
    if btype == "code":
        return min(0.25 * (block["content"].count("\n") + 1) + 0.3, 3.5)
    if btype == "blockquote":
        return 0.7
    if btype == "heading":
        return 0.45
    return 0.3


# ---------------------------------------------------------------------------
# Main conversion logic
# ---------------------------------------------------------------------------
def convert_md_to_pptx(md_path: str, output_path: str, title: str = None):
    """Convert a markdown file to a presentation-quality PowerPoint deck."""
    md_text = Path(md_path).read_text(encoding="utf-8")

    parser = MarkdownParser()
    blocks = parser.parse(md_text)

    planner = PresentationPlanner()
    slide_plan = planner.plan(blocks, title)

    prs = Presentation()
    builder = SlideBuilder(prs, title or "Presentation")

    for planned in slide_plan:
        kind = planned["kind"]
        if kind == "title":
            builder.add_title_slide(planned["title"], planned.get("subtitle", ""),
                                    notes=planned["notes_blocks"] or None)
        elif kind == "section":
            builder.add_section_slide(planned["title"],
                                      notes=planned["notes_blocks"] or None)
        elif kind == "content":
            slide_blocks = planned["slide_blocks"]
            notes_blocks = list(planned["notes_blocks"])

            # Instead of splitting overflow into extra slides, keep ONE slide
            # and move overflow blocks to speaker notes.
            on_slide: list[dict] = []
            est = 0.0
            overflow = False

            for b in slide_blocks:
                h = _estimate_height(b)
                if est + h > 5.0 and on_slide:
                    overflow = True
                    notes_blocks.append(b)
                else:
                    on_slide.append(b)
                    est += h

            builder.add_content_slide(planned["title"], on_slide,
                                      notes_blocks or None)

    prs.save(output_path)
    print(f"✅ Generated: {output_path} ({builder.slide_num} slides)")
    return output_path


# ---------------------------------------------------------------------------
# CLI entry point
# ---------------------------------------------------------------------------
def main():
    ap = argparse.ArgumentParser(
        prog="md2pptx",
        description="Convert Markdown to presentation-quality PowerPoint decks.",
        epilog="Example: python md2pptx.py presentation.md output.pptx --title 'My Talk'",
    )
    ap.add_argument("input", help="Path to the input Markdown file")
    ap.add_argument("output", nargs="?", default=None,
                    help="Path for the output .pptx file (default: same name as input)")
    ap.add_argument("--title", "-t", default=None,
                    help="Override the presentation title")

    args = ap.parse_args()
    input_path = Path(args.input)
    if not input_path.exists():
        print(f"❌ File not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    output_path = args.output or str(input_path.with_suffix(".pptx"))
    convert_md_to_pptx(str(input_path), output_path, args.title)


if __name__ == "__main__":
    main()
